# # app.py - Multi-format Text Extractor

# from flask import Flask, render_template, request, send_file, session, redirect, url_for
# import os, zipfile, datetime
# from PIL import Image, ImageOps
# import pytesseract
# import fitz  # PyMuPDF
# import docx2txt
# import openpyxl
# from pptx import Presentation
# from docx import Document
# from werkzeug.utils import secure_filename

# pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

# app = Flask(__name__)
# app.secret_key = 'secret-key'
# app.config['UPLOAD_FOLDER'] = 'uploads'
# app.config['OUTPUT_FOLDER'] = 'output'

# os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
# os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok=True)

# pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# ALLOWED_EXTENSIONS = {'png', 'jpg', 'jpeg', 'zip', 'pdf', 'docx', 'xlsx', 'pptx'}

# def allowed_file(filename):
#     return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

# def preprocess_image(image_path):
#     img = Image.open(image_path).convert('L')
#     img = ImageOps.invert(img)
#     return img

# def extract_text_from_image(image_path):
#     img = preprocess_image(image_path)
#     return pytesseract.image_to_string(img)

# def extract_text_from_pdf(path):
#     text = ""
#     pdf = fitz.open(path)
#     for page in pdf:
#         text += page.get_text()
#     return text

# def extract_text_from_docx(path):
#     return docx2txt.process(path)

# def extract_text_from_xlsx(path):
#     wb = openpyxl.load_workbook(path)
#     text = ""
#     for sheet in wb.worksheets:
#         for row in sheet.iter_rows():
#             for cell in row:
#                 if cell.value:
#                     text += f"{cell.value}\t"
#             text += "\n"
#     return text

# def extract_text_from_pptx(path):
#     prs = Presentation(path)
#     text = ""
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if hasattr(shape, "text"):
#                 text += shape.text + "\n"
#     return text

# def handle_uploaded_file(filepath, filename):
#     extension = filename.rsplit('.', 1)[1].lower()
#     content = []

#     if extension == 'zip':
#         with zipfile.ZipFile(filepath, 'r') as zip_ref:
#             zip_ref.extractall(app.config['UPLOAD_FOLDER'])
#             for f in zip_ref.namelist():
#                 if f.lower().endswith(('png', 'jpg', 'jpeg')):
#                     image_path = os.path.join(app.config['UPLOAD_FOLDER'], f)
#                     text = extract_text_from_image(image_path)
#                     content.append((f"Session: {f}", text))
#     elif extension in {'png', 'jpg', 'jpeg'}:
#         text = extract_text_from_image(filepath)
#         content.append((f"Session: {filename}", text))
#     elif extension == 'pdf':
#         text = extract_text_from_pdf(filepath)
#         content.append((f"Session: {filename}", text))
#     elif extension == 'docx':
#         text = extract_text_from_docx(filepath)
#         content.append((f"Session: {filename}", text))
#     elif extension == 'xlsx':
#         text = extract_text_from_xlsx(filepath)
#         content.append((f"Session: {filename}", text))
#     elif extension == 'pptx':
#         text = extract_text_from_pptx(filepath)
#         content.append((f"Session: {filename}", text))

#     return content

# @app.route('/')
# def index():
#     return render_template('index.html')

# @app.route('/upload', methods=['POST'])
# def upload():
#     uploaded_file = request.files['file']
#     if not uploaded_file or not allowed_file(uploaded_file.filename):
#         return "Invalid file format."

#     filename = secure_filename(uploaded_file.filename)
#     filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
#     uploaded_file.save(filepath)

#     session['filename_base'] = os.path.splitext(filename)[0]
#     extracted_content = handle_uploaded_file(filepath, filename)
#     session['extracted_text'] = extracted_content
#     return render_template('edit_preview.html', content=extracted_content)

# @app.route('/download', methods=['POST'])
# def download():
#     updated_text = []
#     for key in request.form:
#         updated_text.append((key, request.form[key]))

#     doc = Document()
#     for title, text in updated_text:
#         doc.add_heading(title, level=1)
#         doc.add_paragraph(text)

#     filename_base = session.get('filename_base', 'extracted_text')
#     timestamp = datetime.datetime.now().strftime('%Y%m%d_%H%M%S')
#     output_filename = f"{filename_base}_{timestamp}.docx"
#     output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
#     doc.save(output_path)
#     return send_file(output_path, as_attachment=True)

# # if __name__ == '__main__':
# #    app.run(debug=True, use_reloader=False)








# if __name__ == '__main__':
#     print("🔥 Flask server is starting on http://127.0.0.1:5000")
#  app.run(debug=True, use_reloader=False)



# app.py - Multi-format Text Extractor

from flask import Flask, render_template, request, send_file, session, redirect, url_for
import os, zipfile, datetime
from PIL import Image, ImageOps
import pytesseract
import fitz  # PyMuPDF
import docx2txt
import openpyxl
from pptx import Presentation
from docx import Document
from werkzeug.utils import secure_filename

pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

app = Flask(__name__)
app.secret_key = os.environ.get('SECRET_KEY', 'fallback-secret-key')
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['OUTPUT_FOLDER'] = 'output'

os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok=True)

ALLOWED_EXTENSIONS = {'png', 'jpg', 'jpeg', 'zip', 'pdf', 'docx', 'xlsx', 'pptx'}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def preprocess_image(image_path):
    img = Image.open(image_path).convert('L')         # Convert to grayscale
    img = ImageOps.invert(img)                        # Invert black background to white
    return img

def extract_text_from_image(image_path):
    img = preprocess_image(image_path)
    return pytesseract.image_to_string(img)

def extract_text_from_pdf(path):
    text = ""
    pdf = fitz.open(path)
    for page in pdf:
        text += page.get_text()
    return text

def extract_text_from_docx(path):
    return docx2txt.process(path)

def extract_text_from_xlsx(path):
    wb = openpyxl.load_workbook(path)
    text = ""
    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value:
                    text += f"{cell.value}\t"
            text += "\n"
    return text

def extract_text_from_pptx(path):
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def handle_uploaded_file(filepath, filename):
    extension = filename.rsplit('.', 1)[1].lower()
    content = []

    if extension == 'zip':
        with zipfile.ZipFile(filepath, 'r') as zip_ref:
            zip_ref.extractall(app.config['UPLOAD_FOLDER'])
            for f in zip_ref.namelist():
                if f.lower().endswith(('png', 'jpg', 'jpeg')):
                    image_path = os.path.join(app.config['UPLOAD_FOLDER'], f)
                    if os.path.exists(image_path):  # Check existence
                        text = extract_text_from_image(image_path)
                        content.append((f"Session: {f}", text))
    elif extension in {'png', 'jpg', 'jpeg'}:
        text = extract_text_from_image(filepath)
        content.append((f"Session: {filename}", text))
    elif extension == 'pdf':
        text = extract_text_from_pdf(filepath)
        content.append((f"Session: {filename}", text))
    elif extension == 'docx':
        text = extract_text_from_docx(filepath)
        content.append((f"Session: {filename}", text))
    elif extension == 'xlsx':
        text = extract_text_from_xlsx(filepath)
        content.append((f"Session: {filename}", text))
    elif extension == 'pptx':
        text = extract_text_from_pptx(filepath)
        content.append((f"Session: {filename}", text))

    return content

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload():
    uploaded_file = request.files['file']
    if not uploaded_file or not allowed_file(uploaded_file.filename):
        return "Invalid file format."

    filename = secure_filename(uploaded_file.filename)
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    uploaded_file.save(filepath)

    session['filename_base'] = os.path.splitext(filename)[0]
    extracted_content = handle_uploaded_file(filepath, filename)
    session['extracted_text'] = extracted_content
    return render_template('edit_preview.html', content=extracted_content)

@app.route('/download', methods=['POST'])
def download():
    updated_text = []
    for key in request.form:
        updated_text.append((key, request.form[key]))

    doc = Document()
    for title, text in updated_text:
        doc.add_heading(title, level=1)
        doc.add_paragraph(text)

    filename_base = session.get('filename_base', 'extracted_text')
    timestamp = datetime.datetime.now().strftime('%Y%m%d_%H%M%S')
    output_filename = f"{filename_base}_{timestamp}.docx"
    output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
    doc.save(output_path)
    return send_file(output_path, as_attachment=True)

# if __name__ == '__main__':
#     print("🔥 Flask server is starting on http://127.0.0.1:5000")
#     app.run(debug=True, use_reloader=False)

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port)
